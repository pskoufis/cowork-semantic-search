"""Per-format text extraction from document files.

Spreadsheet formats (.csv, .xlsx, .xlsm, .xls) do not flow through
extract_text — they route via server/spreadsheets.py to a description-based
queue (preview → LLM description → embedded chunk). server.indexer routes
them directly; extract_text rejects them so a stray caller cannot fall into
a raw-row scheme.
"""

import re
from html.parser import HTMLParser
from pathlib import Path

SUPPORTED_EXTENSIONS = {
    ".txt", ".md", ".pdf", ".docx", ".pptx",
    # Spreadsheets (.csv/.xlsx/.xlsm/.xls) temporarily disabled.
    # .mbox, .msg, and .pst are NOT here: the indexer preprocesses every
    # .mbox into a sibling <stem>_unpacked/ tree of per-message .txt files
    # (plus materialized attachments under attachments/msg-NNNN/), every
    # .msg into a sibling <stem>.txt + attachments/<stem>__* layout, and
    # every .pst into a sibling <stem>_unpacked/ tree of per-message .txt
    # files mirroring the PST folder hierarchy. The normal walk picks
    # those up via their native extensions. See
    # mbox_handling.unpack.ensure_unpacked(),
    # msg_handling.unpack.ensure_unpacked(),
    # pst_handling.unpack.ensure_unpacked(), and the preprocessing pass
    # in server.indexer.index_folder.
}


def extract_text(file_path: Path) -> list[dict]:
    """Extract text from a file, returning list of {text, metadata} dicts.

    Metadata may include page_number (PDF) or slide_number (PPTX).
    Spreadsheets are not handled here — see the module docstring.
    """
    suffix = file_path.suffix.lower()

    match suffix:
        case ".txt" | ".md":
            text = file_path.read_text(encoding="utf-8", errors="replace")
            return [{"text": text, "metadata": {}}]
        case ".pdf":
            return _extract_pdf(file_path)
        case ".docx":
            return _extract_docx(file_path)
        case ".pptx":
            return _extract_pptx(file_path)
        case _:
            raise ValueError(
                f"Unsupported file type: {suffix}. "
                f"Supported via extract_text: .txt, .md, .pdf, .docx, .pptx. "
                f"Spreadsheets (.csv/.xlsx/.xlsm/.xls) route through "
                f"server.spreadsheets. .mbox, .msg, and .pst are "
                f"preprocessed into sibling <stem>_unpacked/ trees by "
                f"server.indexer."
            )


def _extract_pdf(file_path: Path) -> list[dict]:
    import pymupdf

    doc = pymupdf.open(str(file_path))
    parts = []
    for page_num, page in enumerate(doc, start=1):
        text = page.get_text()
        parts.append({"text": text.strip(), "metadata": {"page_number": page_num}})
    doc.close()
    return parts


def _extract_docx(file_path: Path) -> list[dict]:
    from docx import Document

    doc = Document(str(file_path))
    text = "\n\n".join(p.text for p in doc.paragraphs)
    return [{"text": text, "metadata": {}}]


def _extract_pptx(file_path: Path) -> list[dict]:
    from pptx import Presentation

    prs = Presentation(str(file_path))
    parts = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        parts.append({
            "text": "\n".join(texts),
            "metadata": {"slide_number": slide_num},
        })
    return parts


# --- HTML / RTF body helpers ------------------------------------------------
#
# Used by msg_handling and pst_handling to reduce an HTML or RTF body to
# plain text when no plain-text body is available. Kept here so both
# packages share one implementation.


class _HTMLTextExtractor(HTMLParser):
    """Collect visible text from an HTML document, dropping script/style."""

    def __init__(self):
        super().__init__(convert_charrefs=True)
        self._chunks: list[str] = []
        self._suppress = 0

    def handle_starttag(self, tag, attrs):
        if tag in ("script", "style"):
            self._suppress += 1

    def handle_endtag(self, tag):
        if tag in ("script", "style") and self._suppress:
            self._suppress -= 1

    def handle_data(self, data):
        if not self._suppress:
            self._chunks.append(data)

    def get_text(self) -> str:
        return "".join(self._chunks)


def _strip_html(html_text: str) -> str:
    """Reduce an HTML body to readable plain text."""
    extractor = _HTMLTextExtractor()
    try:
        extractor.feed(html_text)
        text = extractor.get_text()
    except Exception:
        text = re.sub(r"<[^>]+>", " ", html_text)  # crude fallback
    text = re.sub(r"[ \t\r\f\v]+", " ", text)
    text = re.sub(r" *\n *", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _strip_rtf(rtf_text: str) -> str:
    """Best-effort plain text from an RTF body.

    RTF is only the last-resort body fallback (almost all mail has a plain-text
    or HTML body), so this does a crude control-word strip, not a full parse.
    """
    text = re.sub(r"\\par[d]?\b", "\n", rtf_text)
    text = re.sub(r"\\(line|tab)\b", " ", text)
    text = re.sub(r"\\'[0-9a-fA-F]{2}", "", text)
    text = re.sub(r"\\[a-zA-Z]+-?\d*\s?", "", text)
    text = text.replace("{", "").replace("}", "")
    return re.sub(r"\n{3,}", "\n\n", text).strip()
